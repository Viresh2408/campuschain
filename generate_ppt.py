
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ─── Color palette (matching Yellow/White minimalist style) ─────────────────
YELLOW   = RGBColor(0xF5, 0xC5, 0x18)  # accent yellow
BLACK    = RGBColor(0x1A, 0x1A, 0x2E)  # dark background text
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0x55, 0x55, 0x55)
LGRAY    = RGBColor(0xF7, 0xF7, 0xF7)
DARKGRAY = RGBColor(0x33, 0x33, 0x33)

W = Inches(13.33)   # widescreen width
H = Inches(7.5)     # widescreen height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank

# ─── Helper functions ────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    fill = shape.fill
    if fill_rgb:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()
    line = shape.line
    if line_rgb:
        line.color.rgb = line_rgb
        if line_width:
            line.width = Pt(line_width)
    else:
        line.fill.background()
    return shape

def add_textbox(slide, text, l, t, w, h,
                font_size=14, bold=False, color=BLACK,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox

def add_multiline_textbox(slide, lines, l, t, w, h,
                           base_size=12, base_color=DARKGRAY,
                           base_bold=False, line_spacing=1.15):
    """lines = list of (text, size, bold, color, italic)"""
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, size, bold, color, italic = item, base_size, base_bold, base_color, False
        else:
            text = item[0]
            size  = item[1] if len(item) > 1 else base_size
            bold  = item[2] if len(item) > 2 else base_bold
            color = item[3] if len(item) > 3 else base_color
            italic= item[4] if len(item) > 4 else False

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txBox

def add_bullet_section(slide, heading, bullets, l, t, w, h,
                        bullet_size=11.5, heading_size=13):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    # heading
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = heading
    run.font.size = Pt(heading_size)
    run.font.bold = True
    run.font.color.rgb = YELLOW
    run.font.name = "Calibri"
    # bullets
    for b in bullets:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"  • {b}"
        run.font.size = Pt(bullet_size)
        run.font.bold = False
        run.font.color.rgb = DARKGRAY
        run.font.name = "Calibri"

def yellow_accent_bar(slide):
    """Left yellow vertical bar decoration."""
    add_rect(slide, 0, 0, 0.12, 7.5, fill_rgb=YELLOW)

def slide_header(slide, title, subtitle=None):
    """Standard slide header with yellow underline."""
    yellow_accent_bar(slide)
    add_textbox(slide, title, 0.4, 0.25, 12.4, 0.65,
                font_size=26, bold=True, color=BLACK, align=PP_ALIGN.LEFT)
    add_rect(slide, 0.4, 0.95, 2.0, 0.04, fill_rgb=YELLOW)
    if subtitle:
        add_textbox(slide, subtitle, 0.4, 1.0, 12.4, 0.4,
                    font_size=12, bold=False, color=GRAY, align=PP_ALIGN.LEFT, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Cover
# ════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)

# Background
add_rect(s1, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
# Left yellow stripe
add_rect(s1, 0, 0, 0.5, 7.5, fill_rgb=YELLOW)
# Bottom decoration bar
add_rect(s1, 0, 7.1, 13.33, 0.4, fill_rgb=YELLOW)

# Project title
add_textbox(s1, "CampusChain", 0.9, 1.5, 12.0, 1.1,
            font_size=52, bold=True, color=BLACK, align=PP_ALIGN.LEFT)
add_textbox(s1, "Smart Campus Financial Ecosystem", 0.9, 2.6, 12.0, 0.6,
            font_size=22, bold=False, color=GRAY, align=PP_ALIGN.LEFT)
add_rect(s1, 0.9, 3.35, 3.5, 0.06, fill_rgb=YELLOW)

# Metadata labels
fields = [
    ("PROBLEM STATEMENT:",  "Unified Blockchain-Powered Campus Financial OS for Institutions",       3.7, 14),
    ("COLLEGE NAME:",        "A.P. Shah Institute of Technology, Thane",                             4.25, 13),
    ("TEAM NAME:",           "ViTi",                                                                   4.8, 13),
    ("TEAM MEMBERS:",        "Viresh Nair  |  Tanay Bhatkar  |  [Team ViTi]",                         5.35, 13),
]
for label, val, y, fsz in fields:
    add_textbox(s1, label, 0.9, y, 3.2, 0.45,
                font_size=9, bold=True, color=GRAY, align=PP_ALIGN.LEFT)
    add_textbox(s1, val, 4.1, y, 8.5, 0.45,
                font_size=fsz, bold=True, color=BLACK, align=PP_ALIGN.LEFT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
add_rect(s2, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
slide_header(s2, "Problem Statement",
             "Identifying gaps in how campus finances are currently managed")

content_blocks = [
    ("Clearly Defined Problem", [
        "Campus financial systems operate in fragmented silos — fee payments,",
        "vendor transactions, event ticketing, and peer-to-peer transfers are all",
        "handled by separate, disconnected tools with no unified oversight."
    ]),
    ("Domain & Context", [
        "Higher Education | FinTech | Blockchain Technology",
        "Affects students, faculty, vendors, and administrative staff daily."
    ]),
    ("Key Challenges", [
        "Cash dependency on campus — no digital-first unified payment layer.",
        "No transparent audit trail for inter-campus financial transactions.",
        "Manual fee collection is slow, error-prone, and hard to reconcile.",
        "Students have no visibility into spending patterns or balances."
    ]),
    ("Existing Inefficiencies & Gaps", [
        "Current UPI/bank solutions lack campus-specific role-based controls.",
        "No mechanism for admin-controlled token minting or disbursement.",
        "Event ticketing involves paper-based or third-party tools (Insider, etc.).",
        "Zero fraud detection or anomaly alerting in any existing campus system."
    ]),
    ("Impact on Stakeholders", [
        "Students: financial opacity, no instant transfer capability, cash hassles.",
        "Admins: no real-time dashboards, manual reconciliation burden.",
        "Vendors: delayed settlements, no digital transaction records.",
        "Institution: compliance risk, no audit trail, fragmented reporting."
    ]),
]

y_cursor = 1.45
for heading, bullets in content_blocks:
    add_bullet_section(s2, heading, bullets, 0.4, y_cursor, 12.5, 1.1,
                       bullet_size=10.5, heading_size=12)
    y_cursor += len(bullets) * 0.22 + 0.5
    if y_cursor > 7.0:
        break

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Proposed Solution
# ════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
add_rect(s3, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
slide_header(s3, "Proposed Solution",
             "A unified, blockchain-backed campus financial operating system")

add_textbox(s3,
    "CampusChain is a full-stack, role-aware financial platform that runs on top of a "
    "Supabase-powered PostgreSQL backend with a blockchain-inspired immutable ledger. "
    "It allows admins to mint CampusCoin (CC), students to pay fees, transfer funds P2P, "
    "buy event tickets, and vendors to accept digital QR payments — all in one dashboard.",
    0.4, 1.35, 12.5, 0.9, font_size=11.5, color=DARKGRAY)

# System design diagram (text-based)
add_textbox(s3, "High-Level System Architecture", 0.4, 2.35, 12.5, 0.4,
            font_size=13, bold=True, color=BLACK)
add_rect(s3, 0.4, 2.05, 12.5, 0.04, fill_rgb=YELLOW)

# Draw architecture boxes
boxes = [
    (0.5,  2.85, 2.8, 1.05, "Frontend\n(Next.js + React)", LGRAY),
    (3.7,  2.85, 2.8, 1.05, "API Layer\n(Next.js Routes)", LGRAY),
    (6.9,  2.85, 2.8, 1.05, "Supabase\n(PostgreSQL + Auth)", LGRAY),
    (10.1, 2.85, 2.8, 1.05, "Blockchain Ledger\n(SHA-256 Blocks)", LGRAY),
]
for bx, by, bw, bh, label, bg in boxes:
    add_rect(s3, bx, by, bw, bh, fill_rgb=bg, line_rgb=YELLOW, line_width=1.2)
    add_textbox(s3, label, bx+0.1, by+0.15, bw-0.2, bh-0.2,
                font_size=10.5, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

# Arrows (text →)
for ax in [3.35, 6.55, 9.75]:
    add_textbox(s3, "→", ax, 3.15, 0.4, 0.45,
                font_size=18, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

# Workflow explanation
add_textbox(s3, "Component Interaction Flow", 0.4, 4.1, 12.5, 0.4,
            font_size=13, bold=True, color=BLACK)
add_rect(s3, 0.4, 3.85, 12.5, 0.04, fill_rgb=YELLOW)

workflow_lines = [
    ("1. User authenticates → JWT issued → stored in localStorage.", 11, False, DARKGRAY),
    ("2. All API calls carry the Bearer token → server validates via jsonwebtoken.", 11, False, DARKGRAY),
    ("3. Transactions are written to the Ledger table + a SHA-256 Block is mined.", 11, False, DARKGRAY),
    ("4. Supabase Realtime broadcasts balance changes → dashboard updates live.", 11, False, DARKGRAY),
    ("5. Fraud engine checks rules post-transaction → flags inserted into fraud_alerts.", 11, False, DARKGRAY),
    ("6. EmailJS sends transaction receipt to both sender and receiver.", 11, False, DARKGRAY),
]
add_multiline_textbox(s3, workflow_lines, 0.4, 4.25, 12.5, 2.8)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Approach / Methodology
# ════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
add_rect(s4, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
slide_header(s4, "Approach / Methodology",
             "Step-by-step engineering strategy adopted to build CampusChain")

steps = [
    ("Step 1 — Database-First Design",
     "Modeled all entities (users, ledger, blocks, events, tickets, fraud_alerts) in "
     "PostgreSQL via Supabase. Row-Level Security policies lock every resource to the "
     "correct role (student / admin / vendor)."),

    ("Step 2 — Blockchain-Inspired Immutability",
     "Every financial transaction generates a unique SHA-256 tx_id and a corresponding "
     "Block record (block_number, prev_hash, merkle_root). This creates an append-only "
     "audit chain that cannot be altered without breaking hash continuity."),

    ("Step 3 — JWT Role-Based Auth",
     "Custom /api/auth/login endpoint issues JWTs with embedded {id, role}. All API "
     "routes validate the token server-side using jsonwebtoken, enforcing least-privilege "
     "access control per user type."),

    ("Step 4 — Modular API Architecture",
     "Each major domain (wallet, ledger, events, fraud, analytics) has its own /api/* "
     "route. This separation of concerns simplifies testing, debugging, and future "
     "microservice extraction."),

    ("Step 5 — Rule-Based Fraud Detection",
     "A fraud engine evaluates every outbound transaction against 4 heuristic rules: "
     "(a) 3× average transaction amount, (b) >5 txns in 10 minutes, (c) first-time large "
     "transfer to unknown recipient, (d) suspicious round amounts. Flags are stored in "
     "fraud_alerts and surfaced on the admin dashboard."),

    ("Step 6 — Real-Time UX",
     "Supabase Realtime subscriptions push balance changes and incoming transfers to "
     "connected clients instantly. Toast notifications alert users without a page refresh, "
     "keeping the experience fast and interactive."),
]

y = 1.42
for i, (title, body) in enumerate(steps):
    col = 0.4 if i % 2 == 0 else 6.8
    row = i // 2
    ty = 1.42 + row * 1.95
    add_rect(s4, col, ty, 6.1, 1.75, fill_rgb=LGRAY, line_rgb=YELLOW, line_width=0.8)
    add_textbox(s4, title, col+0.15, ty+0.1, 5.8, 0.38,
                font_size=11, bold=True, color=BLACK)
    add_textbox(s4, body, col+0.15, ty+0.48, 5.8, 1.2,
                font_size=9.5, bold=False, color=DARKGRAY, wrap=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Tech Stack
# ════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
add_rect(s5, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
slide_header(s5, "Tech Stack (Proposed & Implemented)",
             "Every technology was chosen for its real-world production readiness")

tech_cols = [
    ("Frontend", [
        ("Next.js 16",          "React-based framework with server-side rendering & file-based routing (Pages Router)"),
        ("React 19",            "Component-driven UI with hooks for state, effects, and real-time updates"),
        ("Tailwind CSS 4",      "Utility-first CSS framework for a responsive, dark-mode-ready design system"),
        ("Lucide React Icons",  "Lightweight, consistent icon set across all UI components"),
        ("Recharts",            "Declarative charting library used in the Analytics & Spending pages"),
    ]),
    ("Backend & DB", [
        ("Next.js API Routes",  "Serverless backend endpoints: /api/auth, /api/wallet, /api/events, /api/fraud…"),
        ("Supabase (PostgreSQL)","Cloud-hosted relational DB with Row-Level Security, Realtime subscriptions & admin SDK"),
        ("Supabase Realtime",   "WebSocket-based pub/sub that pushes live balance and ledger updates to clients"),
    ]),
    ("Security & Auth", [
        ("jsonwebtoken (JWT)",  "Stateless authentication — tokens embed user id + role; validated on every request"),
        ("bcryptjs",            "Salted password hashing (cost=10) for all user password storage and verification"),
    ]),
    ("Utilities", [
        ("Node.js crypto",      "SHA-256 hashing for transaction IDs, block hashes, and Merkle root generation"),
        ("QRCode",              "Server-side QR code generation for vendor payment codes and event tickets"),
        ("EmailJS",             "Client-side email SDK for sending transaction receipts to both parties instantly"),
        ("date-fns",            "Lightweight, tree-shakeable date formatting for transaction timestamps"),
        ("react-hot-toast",     "Real-time notification toasts for balance changes, incoming transfers, and errors"),
    ]),
]

col_x    = [0.35, 4.55, 8.75]   # not enough for 4 cols; use 3 + overflow
row_y    = 1.42
col_w    = 4.1

for ci, (cat_title, items) in enumerate(tech_cols):
    cx = col_x[min(ci, 2)]
    cy = row_y if ci < 3 else row_y + 2.8
    if ci == 3:
        cx = 0.35

    add_rect(s5, cx, cy, col_w, 0.38, fill_rgb=YELLOW)
    add_textbox(s5, cat_title, cx+0.12, cy+0.04, col_w-0.2, 0.32,
                font_size=12, bold=True, color=BLACK)

    iy = cy + 0.45
    for tech, role in items:
        add_textbox(s5, tech, cx+0.12, iy, col_w-0.2, 0.25,
                    font_size=10, bold=True, color=BLACK)
        add_textbox(s5, role, cx+0.12, iy+0.24, col_w-0.2, 0.42,
                    font_size=9, bold=False, color=GRAY, wrap=True)
        iy += 0.68
        if iy > 7.1:
            break

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Expected Impact / Innovation
# ════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
add_rect(s6, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
slide_header(s6, "Expected Impact & Innovation",
             "What makes CampusChain unique compared to existing campus payment tools")

impact_items = [
    ("🔗 Blockchain Audit Trail",
     "Unlike any campus ERP, every transaction is linked in an immutable SHA-256 "
     "block chain. Admins can verify the integrity of every payment since the genesis block."),
    ("🛡️ Built-In Fraud Intelligence",
     "A rule-based fraud detection engine runs on every outbound transaction, checking "
     "velocity, amount anomalies, and unknown recipients — a feature absent in all "
     "standard campus payment apps."),
    ("⚡ Real-Time Balance & Notifications",
     "Supabase Realtime delivers instant push updates to wallets. Students see incoming "
     "money appear live with toast alerts — no page refresh required, no delay."),
    ("🎫 NFT-Like Event Ticketing",
     "Event tickets are on-chain tokens with a unique ticket_token UUID, linked to a "
     "transaction ID in the ledger. Tickets cannot be duplicated or forged."),
    ("📊 Spending Analytics Dashboard",
     "Category-level spending breakdown (fees, P2P, vendor, events) with trend charts "
     "powered by Recharts — giving students financial self-awareness not seen before."),
    ("🏆 Gamified Leaderboard",
     "A community leaderboard ranks students by token activity, encouraging responsible "
     "financial participation and engagement with campus digital services."),
]

y_start = 1.42
for i, (title, body) in enumerate(impact_items):
    col = 0.35 if i % 2 == 0 else 6.75
    row = i // 2
    ty = y_start + row * 1.92
    add_rect(s6, col, ty, 6.2, 1.74, fill_rgb=LGRAY, line_rgb=YELLOW, line_width=0.8)
    add_textbox(s6, title, col+0.15, ty+0.1, 5.9, 0.42,
                font_size=11.5, bold=True, color=BLACK)
    add_textbox(s6, body, col+0.15, ty+0.52, 5.9, 1.1,
                font_size=9.8, bold=False, color=DARKGRAY, wrap=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Feasibility & Scalability
# ════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank_layout)
add_rect(s7, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
slide_header(s7, "Feasibility & Scalability",
             "Designed for real-world deployment from day one")

feasibility = [
    ("Technical Feasibility", [
        "Fully open-source stack (Next.js + Supabase) — zero licensing cost.",
        "Deployed on Vercel (frontend) + Supabase Cloud (backend) in < 30 minutes.",
        "All components are stateless and horizontally scalable.",
        "JWT authentication works without sessions — scales to millions of concurrent users.",
    ]),
    ("Operational Feasibility", [
        "Institution admin creates accounts, mints CC, and monitors from one dashboard.",
        "Students onboard via email registration — no app download required (browser-first).",
        "EmailJS receipts keep all parties informed without additional infrastructure.",
        "QR-based vendor payments need only a smartphone camera — no POS hardware.",
    ]),
    ("Scalability Plan", [
        "Supabase supports connection pooling (PgBouncer) and read replicas for high load.",
        "Vercel Edge networking auto-scales API routes globally with zero config.",
        "Blockchain table can be sharded by academic year (block ranges per term).",
        "Multi-campus rollout: each campus gets an isolated Supabase project/schema.",
    ]),
    ("Future Expansion", [
        "Mobile app (React Native) reusing the same API layer — no backend changes.",
        "Smart contract integration (Ethereum/Solana) for on-chain token minting.",
        "AI-powered fraud scoring (ML model replacing rule-based engine).",
        "Integration with national payment rails (UPI, NEFT) for token top-ups.",
        "OpenID Connect / Aadhaar eKYC for verified student identity assurance.",
    ]),
]

y_cur = 1.42
left_cols = [0.35, 6.75]
for i, (sec_title, bullets) in enumerate(feasibility):
    cx = left_cols[i % 2]
    cy = y_cur + (i // 2) * 2.85
    add_rect(s7, cx, cy, 6.15, 0.38, fill_rgb=YELLOW)
    add_textbox(s7, sec_title, cx+0.12, cy+0.04, 5.9, 0.32,
                font_size=12, bold=True, color=BLACK)
    by = cy + 0.44
    for b in bullets:
        add_textbox(s7, f"• {b}", cx+0.12, by, 5.9, 0.38,
                    font_size=9.8, color=DARKGRAY, wrap=True)
        by += 0.52

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Thank You
# ════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank_layout)
add_rect(s8, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
add_rect(s8, 0, 0, 0.5, 7.5, fill_rgb=YELLOW)
add_rect(s8, 0, 7.1, 13.33, 0.4, fill_rgb=YELLOW)

add_textbox(s8, "Thank You", 0.9, 2.0, 12.0, 1.4,
            font_size=64, bold=True, color=BLACK, align=PP_ALIGN.LEFT)
add_rect(s8, 0.9, 3.5, 4.0, 0.07, fill_rgb=YELLOW)
add_textbox(s8, "CampusChain — Smart Campus Financial Ecosystem",
            0.9, 3.7, 12.0, 0.55, font_size=17, bold=False, color=GRAY,
            align=PP_ALIGN.LEFT)
add_textbox(s8, "Team ViTi  |  A.P. Shah Institute of Technology",
            0.9, 4.3, 12.0, 0.45, font_size=13, bold=False, color=GRAY,
            align=PP_ALIGN.LEFT)

add_textbox(s8, "Built with  ●  Next.js  ●  Supabase  ●  Blockchain  ●  React",
            0.9, 5.4, 12.0, 0.45, font_size=11, bold=False, color=YELLOW,
            align=PP_ALIGN.LEFT, italic=True)

# ─── Save ────────────────────────────────────────────────────────────────────
out = "CampusChain_Presentation.pptx"
prs.save(out)
print(f"✅ Saved: {out}")
